import os
import re
import shutil
import traceback
import uuid
import asyncio
from functools import lru_cache

import pandas as pd
import httpx
import torch
from docx import Document as DocxDocument
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from pydantic import BaseModel
from pptx import Presentation
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer

app = FastAPI(title="PRN222 Local RAG AI Server")

SUPPORTED_MODELS = ["bge-m3", "e5", "phobert"]
loaded_models = {}
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
QWEN_ADAPTER_PATH = os.path.join(BASE_DIR, "qwen_lora")


@lru_cache(maxsize=1)
def load_qwen_model():
    """Load and cache the Qwen base model with the local LoRA adapter."""
    import json

    from peft import PeftModel
    from transformers import AutoModelForCausalLM, AutoTokenizer

    if not os.path.isfile(os.path.join(QWEN_ADAPTER_PATH, "adapter_config.json")):
        raise FileNotFoundError(
            f"Khong tim thay Qwen LoRA adapter tai: {QWEN_ADAPTER_PATH}"
        )

    with open(
        os.path.join(QWEN_ADAPTER_PATH, "adapter_config.json"),
        encoding="utf-8",
    ) as config_file:
        adapter_config = json.load(config_file)

    base_model_name = adapter_config["base_model_name_or_path"]
    model = AutoModelForCausalLM.from_pretrained(
        base_model_name,
        device_map="auto" if torch.cuda.is_available() else None,
        torch_dtype="auto",
        low_cpu_mem_usage=True,
    )
    model = PeftModel.from_pretrained(model, QWEN_ADAPTER_PATH)
    tokenizer = AutoTokenizer.from_pretrained(QWEN_ADAPTER_PATH)
    model.eval()
    return model, tokenizer

def get_model(model_name: str) -> SentenceTransformer:
    if model_name not in SUPPORTED_MODELS:
        raise HTTPException(status_code=400, detail="Mo hinh khong duoc ho tro")
    
    if model_name not in loaded_models:
        print(f"Dang tai mo hinh {model_name} vao bo nho RAM...")
        if model_name == "bge-m3":
            loaded_models[model_name] = SentenceTransformer("BAAI/bge-m3")
        elif model_name == "e5":
            loaded_models[model_name] = SentenceTransformer("intfloat/multilingual-e5-base")
        elif model_name == "phobert":
            loaded_models[model_name] = SentenceTransformer("keepitreal/vietnamese-sbert")
        print(f"Tai mo hinh {model_name} hoan tat!")
        
    return loaded_models[model_name]

class EmbedRequest(BaseModel):
    text: str
    model_name: str = "bge-m3"


@app.get("/")
async def root():
    return {"message": "Welcome to PRN222 Local RAG AI Server. Visit /docs for API documentation."}


@app.get("/api/health")
async def health_check():
    return {
        "status": "ok",
        "models": SUPPORTED_MODELS,
        "qwen": {
            "adapter_found": os.path.isfile(
                os.path.join(QWEN_ADAPTER_PATH, "adapter_config.json")
            ),
            "cuda_available": torch.cuda.is_available(),
            "loaded": load_qwen_model.cache_info().currsize > 0,
        },
        "chat_providers": {
            "gpt": bool(os.getenv("GITHUB_MODELS_TOKEN")),
            "gemini": bool(os.getenv("GEMINI_API_KEY")),
            "qwen": os.path.isfile(
                os.path.join(QWEN_ADAPTER_PATH, "adapter_config.json")
            ),
        },
    }


@app.post("/api/embed")
async def generate_embedding(req: EmbedRequest):
    model = get_model(req.model_name)
    vector = await asyncio.to_thread(model.encode, req.text)
    return {"vector": vector.tolist()}


def parse_pdf_to_markdown(file_path: str) -> str:
    reader = PdfReader(file_path)
    markdown_content = []

    for index, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            markdown_content.append(f"## --- TRANG {index + 1} ---\n{text}")

    return "\n\n".join(markdown_content)


def parse_docx_to_markdown(file_path: str) -> str:
    doc = DocxDocument(file_path)
    markdown_content = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        style_name = paragraph.style.name
        if style_name.startswith("Heading 1"):
            markdown_content.append(f"# {text}")
        elif style_name.startswith("Heading 2"):
            markdown_content.append(f"## {text}")
        elif style_name.startswith("Heading 3"):
            markdown_content.append(f"### {text}")
        else:
            markdown_content.append(text)

    return "\n\n".join(markdown_content)


def parse_pptx_to_markdown(file_path: str) -> str:
    presentation = Presentation(file_path)
    markdown_content = []

    for index, slide in enumerate(presentation.slides):
        markdown_content.append(f"## --- SLIDE {index + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                markdown_content.append(shape.text.strip())

    return "\n\n".join(markdown_content)


def parse_xlsx_to_markdown(file_path: str) -> str:
    workbook = pd.ExcelFile(file_path)
    markdown_content = []

    for sheet_name in workbook.sheet_names:
        data_frame = pd.read_excel(workbook, sheet_name=sheet_name)
        markdown_content.append(f"## Bang tinh: {sheet_name}")
        markdown_content.append(data_frame.to_markdown(index=False))

    return "\n\n".join(markdown_content)


def chunk_fixed_size(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    chunks = []
    step = max(chunk_size - overlap, 1)

    for start in range(0, len(text), step):
        chunk = text[start : start + chunk_size].strip()
        if chunk:
            chunks.append(chunk)
        if start + chunk_size >= len(text):
            break

    return chunks


def chunk_by_sentence(text: str, max_chunk_size: int = 800) -> list[str]:
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+|\n+", text) if s.strip()]
    chunks = []
    current_chunk = ""

    for sentence in sentences:
        if len(sentence) > max_chunk_size:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
                current_chunk = ""
            chunks.extend(chunk_fixed_size(sentence, max_chunk_size, 0))
            continue

        separator = " " if current_chunk else ""
        if len(current_chunk) + len(separator) + len(sentence) > max_chunk_size:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = sentence
        else:
            current_chunk += separator + sentence

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def chunk_markdown(markdown_result: str, chunk_strategy: str, chunk_size: int = 500) -> list[str]:
    chunk_strategy = (chunk_strategy or "markdown_header").lower()
    chunks = []

    if chunk_strategy == "markdown_header":
        raw_sections = markdown_result.split("\n#")
        for section in raw_sections:
            if section.strip():
                chunks.append("#" + section if not section.startswith("#") else section)
        return chunks

    if chunk_strategy == "fixed_size":
        return chunk_fixed_size(markdown_result, chunk_size, 0)

    if chunk_strategy == "fixed_size_overlap":
        # Overlap default là 50 hoặc 100
        overlap = min(100, chunk_size // 5)
        return chunk_fixed_size(markdown_result, chunk_size, overlap)

    if chunk_strategy == "sentence":
        return chunk_by_sentence(markdown_result, chunk_size)

    if chunk_strategy != "paragraph":
        raise HTTPException(status_code=400, detail=f"Khong ho tro chunk_strategy {chunk_strategy}")

    paragraphs = markdown_result.split("\n\n")
    current_chunk = ""
    for paragraph in paragraphs:
        if len(current_chunk) + len(paragraph) < chunk_size:
            current_chunk += paragraph + "\n\n"
        else:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = paragraph + "\n\n"

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


@app.post("/api/parse-document")
async def parse_and_process_document(
    file: UploadFile = File(...),
    model_name: str = Form("bge-m3"),
    chunk_strategy: str = Form("markdown_header"),
    chunk_size: int = Form(500),
):
    if model_name not in SUPPORTED_MODELS:
        raise HTTPException(status_code=400, detail="Mo hinh khong duoc ho tro")

    safe_filename = os.path.basename(file.filename or "uploaded_file")
    temp_dir = "temp_files"
    os.makedirs(temp_dir, exist_ok=True)
    temp_filename = f"{uuid.uuid4().hex}_{safe_filename}"
    file_path = os.path.join(temp_dir, temp_filename)

    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        extension = os.path.splitext(safe_filename)[1].lower()
        if extension == ".pdf":
            markdown_result = parse_pdf_to_markdown(file_path)
        elif extension in [".docx", ".doc"]:
            markdown_result = parse_docx_to_markdown(file_path)
        elif extension in [".pptx", ".ppt"]:
            markdown_result = parse_pptx_to_markdown(file_path)
        elif extension in [".xlsx", ".xls"]:
            markdown_result = parse_xlsx_to_markdown(file_path)
        else:
            raise HTTPException(status_code=400, detail=f"Khong ho tro file {extension}")

        chunks = chunk_markdown(markdown_result, chunk_strategy, chunk_size)
        model = get_model(model_name)
        processed_chunks = []
        for index, chunk_text in enumerate(chunks):
            vector = await asyncio.to_thread(model.encode, chunk_text)
            processed_chunks.append({
                "chunk_index": index,
                "content": chunk_text,
                "vector": vector.tolist(),
            })

        return {
            "filename": safe_filename,
            "markdown": markdown_result,
            "total_chunks": len(processed_chunks),
            "chunks": processed_chunks,
        }
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)


RETRYABLE_STATUS_CODES = {429, 503}
RETRY_BACKOFF_SECONDS = [5, 20, 45]


def _parse_retry_delay(response: httpx.Response) -> float | None:
    """Lay thoi gian cho tu Retry-After header hoac retryInfo cua Gemini."""
    retry_after = response.headers.get("retry-after")
    if retry_after:
        try:
            return float(retry_after)
        except ValueError:
            pass
    try:
        details = response.json().get("error", {}).get("details", [])
        for detail in details:
            delay = detail.get("retryDelay")
            if delay:
                return float(delay.rstrip("s"))
    except Exception:
        pass
    return None


async def post_json_with_retry(
    url: str,
    headers: dict,
    payload: dict,
    provider_name: str,
    retryable_status_codes: set[int] = RETRYABLE_STATUS_CODES,
) -> httpx.Response:
    """POST voi retry/backoff cho cac ma loi tam thoi (mac dinh 429, 503)."""
    async with httpx.AsyncClient(timeout=120) as client:
        for attempt, backoff in enumerate([*RETRY_BACKOFF_SECONDS, None]):
            response = await client.post(url, headers=headers, json=payload)
            if response.status_code not in retryable_status_codes or backoff is None:
                return response
            delay = min(_parse_retry_delay(response) or backoff, 60)
            print(
                f"[{provider_name}] HTTP {response.status_code}, "
                f"retry {attempt + 1}/{len(RETRY_BACKOFF_SECONDS)} sau {delay:.0f}s"
            )
            await asyncio.sleep(delay)
    return response


class GatewayChatRequest(BaseModel):
    message: str
    provider: str = "gpt"
    max_new_tokens: int = 256
    temperature: float = 0.0


class GatewayChatResponse(BaseModel):
    answer: str
    provider: str
    model: str
    input_tokens: int = 0
    output_tokens: int = 0


def generate_qwen_answer(
    message: str,
    max_new_tokens: int,
    temperature: float,
) -> tuple[str, int, int]:
    """Run blocking model inference outside FastAPI's event loop."""
    model, tokenizer = load_qwen_model()
    model_inputs = tokenizer.apply_chat_template(
        [{"role": "user", "content": message}],
        tokenize=True,
        add_generation_prompt=True,
        return_tensors="pt",
        return_dict=True,
    )
    model_device = next(model.parameters()).device
    model_inputs = {
        key: value.to(model_device)
        for key, value in model_inputs.items()
    }
    generation_options = {
        "max_new_tokens": max_new_tokens,
        "do_sample": temperature > 0,
    }
    if temperature > 0:
        generation_options["temperature"] = temperature

    output_ids = model.generate(**model_inputs, **generation_options)
    input_length = model_inputs["input_ids"].shape[-1]
    output_length = output_ids.shape[-1] - input_length
    answer = tokenizer.decode(
        output_ids[0][input_length:],
        skip_special_tokens=True,
    ).strip()
    return answer, input_length, output_length


async def generate_gpt_answer(request: GatewayChatRequest) -> GatewayChatResponse:
    api_key = os.getenv("GITHUB_MODELS_TOKEN")
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="GITHUB_MODELS_TOKEN chua duoc cau hinh",
        )

    model = os.getenv("GITHUB_MODELS_MODEL", "gpt-4o-mini")
    base_url = os.getenv(
        "GITHUB_MODELS_BASE_URL",
        "https://models.inference.ai.azure.com",
    ).rstrip("/")
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": request.message}],
        "max_tokens": request.max_new_tokens,
        "temperature": request.temperature,
    }
    headers = {"Authorization": f"Bearer {api_key}"}

    response = await post_json_with_retry(
        f"{base_url}/chat/completions",
        headers,
        payload,
        "gpt",
    )

    if response.is_error:
        raise HTTPException(
            status_code=response.status_code,
            detail=f"GitHub Models error: {response.text}",
        )

    result = response.json()
    usage = result.get("usage", {})
    return GatewayChatResponse(
        answer=result["choices"][0]["message"]["content"].strip(),
        provider="gpt",
        model=model,
        input_tokens=usage.get("prompt_tokens", 0),
        output_tokens=usage.get("completion_tokens", 0),
    )


async def _call_gemini_model(
    model: str,
    api_key: str,
    base_url: str,
    request: GatewayChatRequest,
    retryable_status_codes: set[int] = RETRYABLE_STATUS_CODES,
) -> httpx.Response:
    payload = {
        "contents": [
            {
                "role": "user",
                "parts": [{"text": request.message}],
            }
        ],
        "generationConfig": {
            "temperature": request.temperature,
            "maxOutputTokens": request.max_new_tokens,
        },
    }
    # Gemini 2.5+ la thinking model: tat thinking de token khong bi tieu het
    # vao phan suy nghi (gay response rong khi max tokens nho). Cac model
    # 2.0 tro xuong khong ho tro thinkingConfig nen khong duoc gui.
    if not model.startswith(("gemini-1", "gemini-2.0")):
        payload["generationConfig"]["thinkingConfig"] = {"thinkingBudget": 0}
    headers = {"x-goog-api-key": api_key}

    return await post_json_with_retry(
        f"{base_url}/models/{model}:generateContent",
        headers,
        payload,
        f"gemini:{model}",
        retryable_status_codes,
    )


async def generate_gemini_answer(request: GatewayChatRequest) -> GatewayChatResponse:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="GEMINI_API_KEY chua duoc cau hinh",
        )

    primary_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
    fallback_model = os.getenv("GEMINI_FALLBACK_MODEL", "gemini-3.5-flash")
    base_url = os.getenv(
        "GEMINI_BASE_URL",
        "https://generativelanguage.googleapis.com/v1beta",
    ).rstrip("/")
    has_fallback = bool(fallback_model) and fallback_model != primary_model

    model = primary_model
    if has_fallback:
        # Quota het theo ngay (429) khong the phuc hoi trong vai chuc giay,
        # nen KHONG retry/cho tren model chinh - goi thu 1 lan roi chuyen
        # ngay sang model du phong neu dinh 429. Loi 503 (qua tai tam thoi)
        # thi van dang retry vi co the tu phuc hoi.
        response = await _call_gemini_model(
            model, api_key, base_url, request, retryable_status_codes={503}
        )
        if response.status_code == 429:
            print(f"[gemini] {primary_model} het quota (429), chuyen sang {fallback_model}")
            model = fallback_model
            response = await _call_gemini_model(model, api_key, base_url, request)
    else:
        response = await _call_gemini_model(model, api_key, base_url, request)

    if response.is_error:
        raise HTTPException(
            status_code=response.status_code,
            detail=f"Gemini error ({model}): {response.text}",
        )

    result = response.json()
    usage = result.get("usageMetadata", {})
    candidates = result.get("candidates") or []
    parts = (
        candidates[0].get("content", {}).get("parts") if candidates else None
    ) or []
    answer = "".join(part.get("text", "") for part in parts).strip()
    if not answer:
        finish_reason = (
            candidates[0].get("finishReason", "UNKNOWN")
            if candidates
            else "NO_CANDIDATES"
        )
        raise HTTPException(
            status_code=502,
            detail=f"Gemini tra ve rong (finishReason={finish_reason})",
        )
    return GatewayChatResponse(
        answer=answer,
        provider="gemini",
        model=model,
        input_tokens=usage.get("promptTokenCount", 0),
        output_tokens=usage.get("candidatesTokenCount", 0),
    )


async def route_chat(request: GatewayChatRequest) -> GatewayChatResponse:
    if not request.message.strip():
        raise HTTPException(status_code=400, detail="Message khong duoc de trong")
    if request.max_new_tokens < 1 or request.max_new_tokens > 4096:
        raise HTTPException(
            status_code=400,
            detail="max_new_tokens phai nam trong khoang 1-4096",
        )
    if request.temperature < 0 or request.temperature > 2:
        raise HTTPException(
            status_code=400,
            detail="temperature phai nam trong khoang 0-2",
        )

    provider = request.provider.strip().lower()
    if provider in {"gpt", "chatgpt", "gpt-4o-mini"}:
        return await generate_gpt_answer(request)
    if provider in {"gemini", "gemini-2.5-flash"}:
        return await generate_gemini_answer(request)
    if provider in {"qwen", "qwen-lora"}:
        answer, input_tokens, output_tokens = await asyncio.to_thread(
            generate_qwen_answer,
            request.message,
            request.max_new_tokens,
            request.temperature,
        )
        return GatewayChatResponse(
            answer=answer,
            provider="qwen",
            model="qwen2.5-1.5b-instruct-lora",
            input_tokens=input_tokens,
            output_tokens=output_tokens,
        )

    raise HTTPException(
        status_code=400,
        detail="Provider khong hop le. Ho tro: gpt, gemini, qwen",
    )


@app.post("/api/chat", response_model=GatewayChatResponse)
async def gateway_chat(request: GatewayChatRequest):
    try:
        return await route_chat(request)
    except HTTPException:
        raise
    except Exception as exc:
        error_trace = traceback.format_exc()
        print(error_trace, file=os.sys.stderr)
        error_detail = str(exc).strip() or type(exc).__name__
        raise HTTPException(status_code=500, detail=error_detail) from exc


@app.post("/api/chat/{provider}", response_model=GatewayChatResponse)
async def provider_chat(provider: str, request: GatewayChatRequest):
    request.provider = provider
    return await gateway_chat(request)


@app.post("/api/qwen/chat", response_model=GatewayChatResponse)
@app.post("/qwen-chat")
async def qwen_chat(request: GatewayChatRequest):
    request.provider = "qwen"
    return await gateway_chat(request)
